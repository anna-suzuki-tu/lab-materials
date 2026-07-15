"""
slide_framework.py — v2.1
=========================

SLIDE.md（20260703エンジニアリング協会）のデザインシステムを焼き込んだ
アカデミア発表向け pptx 生成フレームワーク。

デザイン原則（SLIDE.md より）
  - アイスブルー背景 #E7F2F6 のフラット単色。グラデーション・装飾なし
  - 見出し左に幅6px・角丸の Primary 縦バー
  - カードは白・角丸12px・薄い影・1px枠線
  - Accent（ゴールド）は強調したい数値・キーワードのみ
  - ページ番号は右下「N / 25」形式
  - 1スライド1メッセージ

フォント（Mac PowerPoint ネイティブ）
  - 見出し   : Hiragino Sans W8（Noto Sans JP 900 相当）
  - 小見出し : Hiragino Sans W6（700 相当）
  - 本文     : Hiragino Sans W4（400 相当）

パターン（17種）
  cover / closing / section / agenda / key_message
  numbered_cards / icon_list / two_column / before_after / problem_solution
  steps_horizontal / steps_vertical / cards_3 / cards_2x2
  figure_text / table_compare / kpi_3 / profile

共通仕様
  - 本文テキスト中の **テキスト** は Accent 色 + W6 で強調される
  - takeaway= を渡すと下部に全幅のテイクアウェイ帯（ピル型）が付く
    ※ タイトルの言い換えではなく「次のスライドへの橋渡し」か「So What」を書く
  - ページ番号・文字の自動縮小・画像のアスペクト比保持はすべて自動

文字サイズの原則（v2.1 / 2026-07-14 Anna 指示・日英共通）
  - タイトル 28pt ／ 副タイトル（リード文）20pt ／ キーメッセージ 24pt
    ／ 本文 20pt ／ 補足 14pt
  - いずれも「原則」であり、はみ出す場合は自動縮小（_fit_pt / _fit_one_line）
  - 「本文」= カード見出し・箇条書き・ステップ名などの主内容
    「補足」= その下の説明文・キャプション・注記・タグ類

欧文フォント（v2.1 / en・en-ja モード）
  - 見出し・強調 : Avenir Black（Avenir Next Heavy は使わない）
  - 本文・補足   : Avenir Book
"""

import math
import os
import re
import unicodedata

from lxml import etree
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

__version__ = "2.0"

# =========================================================
# Theme tokens（SLIDE.md）
# =========================================================


def _mix(c1, c2, t):
    """c1 を t、c2 を 1-t の割合で混色する。"""
    return RGBColor(
        round(c1[0] * t + c2[0] * (1 - t)),
        round(c1[1] * t + c2[1] * (1 - t)),
        round(c1[2] * t + c2[2] * (1 - t)),
    )


PRIMARY    = RGBColor(0x20, 0x7A, 0x91)   # スチールブルー（縦バー・番号・帯）
SECONDARY  = RGBColor(0x35, 0x6F, 0x8C)   # テールブルー
ACCENT     = RGBColor(0xC6, 0xA4, 0x4A)   # ゴールデンブロンズ（強調数値のみ）
BG         = RGBColor(0xE7, 0xF2, 0xF6)   # アイスブルー
SURFACE    = RGBColor(0xFF, 0xFF, 0xFF)   # カード
BORDER     = RGBColor(0xE5, 0xE7, 0xEB)   # カード枠線
TEXT       = RGBColor(0x1F, 0x29, 0x37)   # ダークチャコール
TEXT_SUB   = RGBColor(0x7D, 0x7E, 0x7F)
TEXT_MUTED = RGBColor(0x5D, 0x5D, 0x5D)
TEAL_LIGHT = RGBColor(0xB0, 0xD4, 0xDE)   # グラフ・矢印の淡色
WHITE      = SURFACE

PRIMARY_TINT   = _mix(PRIMARY, SURFACE, 0.12)    # アイコンチップ背景
SECONDARY_TINT = _mix(SECONDARY, SURFACE, 0.12)
ACCENT_TINT    = _mix(ACCENT, SURFACE, 0.12)
GHOST          = _mix(PRIMARY, BG, 0.07)         # セクション背景の巨大数字
GRAY_CHIP      = RGBColor(0xF0, 0xF0, 0xF0)      # Before ラベル等
ACCENT_DARK    = _mix(ACCENT, TEXT, 0.75)        # ゴールド文字（白地で読める濃度）

CHIP_COLORS = [(PRIMARY_TINT, PRIMARY), (SECONDARY_TINT, SECONDARY),
               (ACCENT_TINT, ACCENT_DARK)]

FONT_HEAVY = "Hiragino Sans W8"   # H1
FONT_BOLD  = "Hiragino Sans W6"   # H2・強調
FONT_BODY  = "Hiragino Sans W4"   # 本文
FONT_LIGHT = "Hiragino Sans W3"   # キャプション

# 言語モード別の欧文フォント（a:latin）。日本語グリフ（a:ea）は常にヒラギノ。
#   ja / ja-en : 欧文もヒラギノ（SLIDE.md「フォントを混在させない」準拠）
#   en / en-ja : 欧文は Avenir（Mac 標準・ヒラギノと好相性）
#   ※ v2.1：Avenir Black + Avenir Book（Heavy 不使用・Anna 指示）
_LATIN_FOR_EN = {
    FONT_HEAVY: "Avenir Black",
    FONT_BOLD:  "Avenir Black",
    FONT_BODY:  "Avenir Book",
    FONT_LIGHT: "Avenir Book",
}

# 文字サイズの原則（pt）。はみ出す場合は自動縮小される
PT_TITLE  = 28   # スライドタイトル
PT_LEAD   = 20   # 副タイトル（リード文・セクション副題）
PT_KEYMSG = 24   # キーメッセージ
PT_BODY   = 20   # 本文（カード見出し・箇条書き・ステップ名）
PT_NOTE   = 14   # 補足（説明文・キャプション・注記）

# =========================================================
# Layout constants（1280x720px 基準 / 96px = 1in）
# =========================================================
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MX  = 0.73                       # 左右余白（70px）
CW  = SLIDE_W_IN - 2 * MX        # コンテンツ幅 11.87in
TITLE_Y = 0.52

IN = Inches
_ACCENT_RE = re.compile(r"\*\*(.+?)\*\*")


# =========================================================
# Text metrics（日本語向け近似計測）
# =========================================================
def _char_w(ch):
    # 欧文は 0.55（Avenir Next 等のやや広い書体でも収まる安全側の係数）
    return 1.0 if unicodedata.east_asian_width(ch) in ("F", "W", "A") else 0.55


def _plain(text):
    return _ACCENT_RE.sub(r"\1", text)


def _est_w_in(text, pt):
    """1行テキストの幅（インチ）を近似する。"""
    return pt / 72.0 * sum(_char_w(c) for c in _plain(text)) * 1.02


def _wrap_lines(text, w_in, pt):
    """幅 w_in に word_wrap したときの行数を近似する。"""
    if not text:
        return 1
    return max(1, math.ceil(_est_w_in(text, pt) / max(w_in, 0.1)))


def _fit_pt(lines, w_in, h_in, base_pt, min_pt, ls=1.3):
    """複数行テキストが h_in に収まる最大フォントサイズを返す。"""
    if isinstance(lines, str):
        lines = [lines]
    pt = base_pt
    while pt > min_pt:
        n = sum(_wrap_lines(l, w_in, pt) for l in lines)
        if n * pt * ls / 72.0 <= h_in:
            return pt
        pt -= 0.5
    return min_pt


def _block_h(lines, w_in, pt, ls=1.3):
    if isinstance(lines, str):
        lines = [lines]
    n = sum(_wrap_lines(l, w_in, pt) for l in lines)
    return n * pt * ls / 72.0


def _fit_one_line(text, w_in, base_pt, min_pt):
    """wrap=False の1行テキストが幅 w_in に収まる最大フォントサイズ。"""
    pt = base_pt
    while pt > min_pt and _est_w_in(text, pt) > w_in:
        pt -= 0.5
    return pt


def _fit_label(text, w_in, h_in, base_pt, min_pt, ls=1.15, one_line_min=13):
    """短いラベル用：まず1行に収まるサイズを試し、無理なら折返しに切替。

    折返し時も最長の単語が途中改行しないよう追加で縮小する
    （日本語のような空白なしテキストには適用しない）。
    """
    pt = _fit_one_line(text, w_in, base_pt, one_line_min)
    if _est_w_in(text, pt) <= w_in:
        return pt
    pt = _fit_pt(text, w_in, h_in, base_pt, min_pt, ls)
    tokens = text.split()
    if len(tokens) > 1:
        widest = max(tokens, key=lambda t: _est_w_in(t, 10))
        while pt > min_pt and _est_w_in(widest, pt) > w_in:
            pt -= 0.5
    return pt


# =========================================================
# Framework
# =========================================================
class SlideFramework:
    def __init__(self, lang="ja"):
        """
        lang : "ja"（日本語のみ）/ "ja-en"（日本語中心・英語補助）
               / "en-ja"（英語中心・日本語補助）/ "en"（英語のみ）
        レイアウトは全モード共通。欧文フォント（a:latin）だけが切り替わり、
        日本語グリフ（a:ea）は常にヒラギノで描画される。
        どのスロットにどちらの言語を書くかは skill（文言設計ルール）側で規定。
        """
        assert lang in ("ja", "ja-en", "en-ja", "en"), f"unknown lang: {lang}"
        self.lang = lang
        self._latin_map = _LATIN_FOR_EN if lang in ("en", "en-ja") else {}
        self.prs = Presentation()
        self.prs.slide_width  = IN(SLIDE_W_IN)
        self.prs.slide_height = IN(SLIDE_H_IN)
        self._page = 0
        self._page_runs = []   # (run, page_no) — save() で "N / total" を確定

    # =====================================================
    # Low-level helpers
    # =====================================================
    def _slide(self, bg=BG):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg
        self._page += 1
        self._page_tag(slide)
        return slide

    def _page_tag(self, slide):
        tb = slide.shapes.add_textbox(IN(SLIDE_W_IN - 2.1), IN(SLIDE_H_IN - 0.42),
                                      IN(2.1 - 0.35), IN(0.28))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = ""
        self._style_run(r, FONT_BODY, 11, TEXT_MUTED)
        self._page_runs.append((r, self._page))

    def _style_run(self, run, font, pt, color, spc=None):
        run.font.size = Pt(pt)
        # a:latin（欧文グリフ）— en 系モードでは Avenir Next に切替
        run.font.name = self._latin_map.get(font, font)
        run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))                # a:ea（日本語グリフ）は常にヒラギノ
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)
        if spc:
            rPr.set("spc", str(int(spc * 100)))  # 字間（pt）

    @staticmethod
    def _soft_shadow(shape, blur_in=0.06, dist_in=0.016, alpha_pct=10):
        """SLIDE.md の `0 1px 4px rgba(0,0,0,0.08)` 相当のソフトシャドウ。"""
        spPr = shape._element.spPr
        eff = spPr.find(qn("a:effectLst"))
        if eff is None:
            eff = etree.SubElement(spPr, qn("a:effectLst"))
        shdw = etree.SubElement(eff, qn("a:outerShdw"), {
            "blurRad": str(int(IN(blur_in))),
            "dist":    str(int(IN(dist_in))),
            "dir":     "5400000",               # 真下
            "rotWithShape": "0",
        })
        clr = etree.SubElement(shdw, qn("a:srgbClr"), {"val": "000000"})
        etree.SubElement(clr, qn("a:alpha"), {"val": str(alpha_pct * 1000)})

    def _rect(self, slide, x, y, w, h, fill, line=None, radius_in=0.0,
              shadow=False, line_w=1.0):
        shape_type = (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius_in > 0
                      else MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        shp = slide.shapes.add_shape(shape_type, IN(x), IN(y), IN(w), IN(h))
        if radius_in > 0:
            shp.adjustments[0] = min(0.5, radius_in / max(min(w, h), 0.01))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        if shadow:
            self._soft_shadow(shp)
        return shp

    def _card(self, slide, x, y, w, h, fill=SURFACE, radius_in=0.125):
        """白カード：角丸12px・1px枠線・薄い影（SLIDE.md Component Style）。"""
        return self._rect(slide, x, y, w, h, fill, line=BORDER,
                          radius_in=radius_in, shadow=True)

    def _circle(self, slide, x, y, d, fill, glyph=None, glyph_pt=None,
                glyph_color=WHITE, glyph_font=FONT_HEAVY):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                     IN(x), IN(y), IN(d), IN(d))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
        shp.shadow.inherit = False
        if glyph is not None:
            self._text(slide, x - 0.1, y, d + 0.2, d, glyph,
                       font=glyph_font, pt=glyph_pt or d * 42,
                       color=glyph_color, align=PP_ALIGN.CENTER,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
        return shp

    # -------- テキスト --------
    def _parse_runs(self, line):
        """`**強調**` を Accent 色 + W6 の run に分解する。"""
        if isinstance(line, list):
            return line
        runs, pos = [], 0
        for m in _ACCENT_RE.finditer(line):
            if m.start() > pos:
                runs.append((line[pos:m.start()], {}))
            runs.append((m.group(1), {"accent": True}))
            pos = m.end()
        if pos < len(line):
            runs.append((line[pos:], {}))
        return runs or [("", {})]

    def _text(self, slide, x, y, w, h, content, font=FONT_BODY, pt=13,
              color=TEXT, align=PP_ALIGN.LEFT,
              valign=MSO_VERTICAL_ANCHOR.TOP, ls=1.3, wrap=True,
              spc=None, accent_color=ACCENT_DARK):
        tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = valign
        lines = [content] if isinstance(content, str) else list(content)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = ls
            for txt, ov in self._parse_runs(line):
                r = p.add_run()
                r.text = txt
                f = ov.get("font") or (FONT_BOLD if ov.get("accent") else font)
                c = ov.get("color") or (accent_color if ov.get("accent") else color)
                self._style_run(r, f, ov.get("pt", pt), c, spc=spc)
        return tb

    def _dot_rows(self, slide, x, y, w, items, pt=13, color=TEXT,
                  dot_color=PRIMARY, ls=1.35, gap=0.14, max_y=None,
                  spread_h=None):
        """Primary 色の丸マーカー付き箇条書き（本物のドットを描画）。

        spread_h を渡すと、その高さの中に項目を縦に均等配分する
        （項目が少ないときに上に詰まって下が空くのを防ぐ）。
        """
        if spread_h and len(items) > 1:
            total = sum(_block_h(it, w - 0.26, pt, ls) for it in items)
            gap = min(0.60, max(gap, (spread_h - total) / (len(items) - 1)))
        cy = y
        for item in items:
            if max_y:
                avail = max_y - cy
                pt_i = _fit_pt(item, w - 0.26, max(avail, 0.3), pt, 10, ls)
            else:
                pt_i = pt
            self._circle(slide, x, cy + pt_i / 72 * ls / 2 - 0.035, 0.09,
                         dot_color)
            self._text(slide, x + 0.26, cy, w - 0.26, 1.0, item,
                       pt=pt_i, color=color, ls=ls)
            cy += _block_h(item, w - 0.26, pt_i, ls) + gap
        return cy

    def _icon_chip(self, slide, x, y, size, icon=None, tint=PRIMARY_TINT,
                   color=PRIMARY):
        """角丸のアイコンチップ。icon は画像パス or 1〜2文字の文字列。"""
        self._rect(slide, x, y, size, size, tint, radius_in=size * 0.25)
        if icon is None:
            return
        if isinstance(icon, str) and not os.path.exists(icon):
            pt = size * (32 if len(_plain(icon)) <= 2 else 20)
            self._text(slide, x, y + 0.01, size, size, icon,
                       font=FONT_HEAVY, pt=pt, color=color,
                       align=PP_ALIGN.CENTER,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
        else:
            pad = size * 0.20
            self._picture_fit(slide, icon, x + pad, y + pad,
                              size - 2 * pad, size - 2 * pad)

    def _picture_fit(self, slide, path, x, y, max_w, max_h,
                     halign="center", valign="middle"):
        """アスペクト比を保持して max_w × max_h 内に画像を配置する。"""
        with Image.open(path) as im:
            iw, ih = im.size
        scale = min(max_w / iw, max_h / ih)
        w, h = iw * scale, ih * scale
        px = x + {"left": 0, "center": (max_w - w) / 2,
                  "right": max_w - w}[halign]
        py = y + {"top": 0, "middle": (max_h - h) / 2,
                  "bottom": max_h - h}[valign]
        return slide.shapes.add_picture(path, IN(px), IN(py), IN(w), IN(h))

    def _chevron(self, slide, x, y, w=0.42, h=0.55, color=TEAL_LIGHT):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON,
                                     IN(x), IN(y), IN(w), IN(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _down_arrow(self, slide, x, y, w=0.34, h=0.26, color=TEAL_LIGHT):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                                     IN(x), IN(y), IN(w), IN(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    # =====================================================
    # 共通フレーム（縦バー見出し + リード文）
    # =====================================================
    def _frame(self, slide, title, lead=None):
        """見出しエリアを描き、コンテンツ開始 Y（インチ）を返す。"""
        w_avail = CW - 0.25
        pt, n_lines = PT_TITLE, 1
        while pt > 23 and _est_w_in(title, pt) > w_avail:
            pt -= 1
        if _est_w_in(title, pt) > w_avail:          # 2行タイトルに切替
            n_lines = 2
            pt = _fit_pt(title, w_avail, 1.0, 24, 20, 1.15)
        th = 0.52 if n_lines == 1 else 0.95
        bar_h = 0.38 if n_lines == 1 else 0.72
        self._rect(slide, MX, TITLE_Y + (th - bar_h) / 2, 0.063, bar_h,
                   PRIMARY, radius_in=0.031)
        self._text(slide, MX + 0.22, TITLE_Y, CW - 0.22, th, title,
                   font=FONT_HEAVY, pt=pt, color=TEXT,
                   valign=MSO_VERTICAL_ANCHOR.MIDDLE, ls=1.12)
        y = TITLE_Y + th + 0.06
        if lead:
            lpt = _fit_one_line(lead, CW - 0.22, PT_LEAD, 11)
            self._text(slide, MX + 0.22, y, CW - 0.22, 0.38, lead,
                       pt=lpt, color=TEXT_MUTED, wrap=False)
            y += 0.56
        else:
            y += 0.14
        return y

    def _takeaway(self, slide, text):
        """下部テイクアウェイ帯（ピル型・Primary）。"""
        if not text:
            return
        pt = _fit_pt(text, 10.4, 0.32, 15, 11.5, 1.0)
        w = min(11.2, max(4.0, _est_w_in(text, pt) + 1.2))
        x = (SLIDE_W_IN - w) / 2
        band = self._rect(slide, x, 6.68, w, 0.54, PRIMARY, radius_in=0.27,
                          shadow=True)
        band.adjustments[0] = 0.5
        self._text(slide, x, 6.68, w, 0.54, text, font=FONT_BOLD, pt=pt,
                   color=WHITE, align=PP_ALIGN.CENTER,
                   valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False,
                   accent_color=_mix(ACCENT, WHITE, 0.55))

    @staticmethod
    def _cb(takeaway):
        """コンテンツ下端 Y。takeaway 帯の有無で変わる。"""
        return 6.52 if takeaway else 6.98

    # =====================================================
    # 1. cover — 表紙（cover-title-center）
    # =====================================================
    def cover(self, title, subtitle=None, event=None, affiliation=None,
              name=None, date=None):
        """
        title    : str または 複数行 list[str]
        subtitle : タイトル下のサブタイトル行（グレー・やや小さく）
        event    : 最上部の会合名（字間を空けた小さめ表示）
        """
        slide = self._slide()
        lines = [title] if isinstance(title, str) else list(title)

        if event:
            self._text(slide, 0, 1.10, SLIDE_W_IN, 0.35, event,
                       font=FONT_BOLD, pt=12, color=TEXT_MUTED,
                       align=PP_ALIGN.CENTER, spc=3.0, wrap=False)

        # タイトルブロック（縦バー + 左揃えテキスト）を水平センタリング
        # wrap=False で描画するため、必ず1行で収まるサイズまで縮小する
        t_pt = 38
        while t_pt > 22 and max(_est_w_in(l, t_pt) for l in lines) > 11.3:
            t_pt -= 1
        s_pt = t_pt * 0.62
        block_h = len(lines) * t_pt * 1.25 / 72 + \
            (s_pt * 1.35 / 72 if subtitle else 0)
        text_w = max(_est_w_in(l, t_pt) for l in lines)
        if subtitle:
            text_w = max(text_w, _est_w_in(subtitle, s_pt))
        bx = max(MX, (SLIDE_W_IN - (0.083 + 0.30 + text_w)) / 2)
        by = 2.55 - block_h / 2
        self._rect(slide, bx, by, 0.083, block_h, PRIMARY, radius_in=0.042)
        content = [[(l, {"pt": t_pt})] for l in lines]
        if subtitle:
            content.append([(subtitle, {"pt": s_pt, "color": TEXT_SUB,
                                        "font": FONT_BOLD})])
        self._text(slide, bx + 0.30, by - 0.05, text_w + 0.6, block_h + 0.1,
                   content, font=FONT_HEAVY, pt=t_pt, color=TEXT,
                   valign=MSO_VERTICAL_ANCHOR.MIDDLE, ls=1.25, wrap=False)

        dy = by + block_h + 0.45
        self._rect(slide, (SLIDE_W_IN - 3.1) / 2, dy, 3.1, 0.032, PRIMARY,
                   radius_in=0.016)
        y = dy + 0.35
        if affiliation:
            self._text(slide, 0, y, SLIDE_W_IN, 0.32, affiliation,
                       pt=13.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
                       wrap=False)
            y += 0.40
        if name:
            self._text(slide, 0, y, SLIDE_W_IN, 0.38, name, font=FONT_BOLD,
                       pt=16.5, color=TEXT, align=PP_ALIGN.CENTER, wrap=False)
            y += 0.46
        if date:
            self._text(slide, 0, y, SLIDE_W_IN, 0.28, date, pt=10.5,
                       color=TEXT_MUTED, align=PP_ALIGN.CENTER, wrap=False)

    # =====================================================
    # 2. section — 章区切り（section-divider・背景に薄い巨大数字）
    # =====================================================
    def section(self, number, title, subtitle=None):
        slide = self._slide()
        self._text(slide, 5.6, 1.1, 7.0, 5.3, f"{number:02d}",
                   font=FONT_HEAVY, pt=270, color=GHOST,
                   align=PP_ALIGN.RIGHT, valign=MSO_VERTICAL_ANCHOR.MIDDLE,
                   wrap=False, ls=1.0)
        block_top = 2.45 if subtitle else 2.75
        t_pt = _fit_pt(title, 10.5, 1.0, 46, 32, 1.0)
        t_h = t_pt / 72 * 1.15
        bar_h = 0.42 + t_h + (0.85 if subtitle else 0.05)
        self._rect(slide, MX, block_top, 0.083, bar_h, PRIMARY,
                   radius_in=0.042)
        x = MX + 0.38
        self._text(slide, x, block_top + 0.02, 9.0, 0.30,
                   f"SECTION {number:02d}", font=FONT_BOLD, pt=11,
                   color=TEXT_MUTED, spc=2.8, wrap=False)
        self._text(slide, x, block_top + 0.42, 10.8, 1.05, title,
                   font=FONT_HEAVY, pt=t_pt, color=TEXT, wrap=False)
        if subtitle:
            spt = _fit_pt(subtitle, 9.6, 0.9, PT_LEAD, 13, 1.5)
            self._text(slide, x, block_top + 0.42 + t_h + 0.22,
                       9.6, 0.9, subtitle, pt=spt, color=TEXT_SUB, ls=1.5)

    # =====================================================
    # 3. agenda — 目次（numbered-toc + 本日の問い）
    # =====================================================
    def agenda(self, items, question=None, title="本日の流れ", lead=None,
               question_label="本日の問い"):
        """
        items    : list[str] または list[(str, str)]（項目, 補足）
        question : 「本日の問い」。指定すると下部にゴールドの問いカードが付く
        question_label : 問いカードのラベル（英語版では "TODAY'S QUESTION" 等）
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = 6.98
        q_h = 0.95 if question else 0.0
        avail = cb - ct - (q_h + 0.35 if question else 0)
        n = len(items)
        row_h = min(0.72, avail / n - 0.14)
        step = min(1.15, avail / n)            # 縦に均等配分
        y0 = ct + (avail - step * n) / 2
        for i, item in enumerate(items):
            text, sub = item if isinstance(item, tuple) else (item, None)
            y = y0 + i * step + (step - row_h) / 2
            cy = y + row_h / 2
            self._circle(slide, MX + 0.05, cy - 0.21, 0.42, PRIMARY,
                         glyph=str(i + 1), glyph_pt=15)
            self._text(slide, MX + 0.72, y, 7.6, row_h, text,
                       font=FONT_BOLD,
                       pt=_fit_one_line(text, 7.6, PT_BODY, 13),
                       color=TEXT,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
            if sub:
                sub_w = SLIDE_W_IN - MX - 8.6
                self._text(slide, 8.6, y, sub_w, row_h, sub,
                           pt=_fit_one_line(sub, sub_w, PT_NOTE, 10),
                           color=TEXT_MUTED,
                           valign=MSO_VERTICAL_ANCHOR.MIDDLE,
                           align=PP_ALIGN.RIGHT, wrap=False)
        if question:
            qy = cb - q_h
            self._card(slide, MX, qy, CW, q_h, fill=ACCENT_TINT)
            self._rect(slide, MX, qy, 0.07, q_h, ACCENT, radius_in=0.035)
            self._text(slide, MX + 0.35, qy + 0.13, 3.5, 0.26,
                       question_label, font=FONT_BOLD, pt=11.5,
                       color=ACCENT_DARK, wrap=False)
            self._text(slide, MX + 0.35, qy + 0.40, CW - 0.7, q_h - 0.5,
                       question, font=FONT_BOLD,
                       pt=_fit_one_line(question, CW - 0.7, PT_BODY, 13),
                       color=TEXT, wrap=False)

    # =====================================================
    # 4. key_message — 強調1枚（key-message-single）
    # =====================================================
    def key_message(self, message, supplement=None):
        """発表の転換点・最重要メッセージ用。1文を中央に大きく。"""
        slide = self._slide()
        self._rect(slide, (SLIDE_W_IN - 0.9) / 2, 2.30, 0.9, 0.028,
                   TEAL_LIGHT, radius_in=0.014)
        pt = _fit_pt(message, 10.6, 1.7, PT_KEYMSG, 18, 1.35)
        self._text(slide, (SLIDE_W_IN - 10.6) / 2, 2.55, 10.6, 1.75, message,
                   font=FONT_HEAVY, pt=pt, color=TEXT, align=PP_ALIGN.CENTER,
                   valign=MSO_VERTICAL_ANCHOR.MIDDLE, ls=1.4)
        self._rect(slide, (SLIDE_W_IN - 0.9) / 2, 4.55, 0.9, 0.028,
                   TEAL_LIGHT, radius_in=0.014)
        if supplement:
            spt = _fit_pt(supplement, 9.0, 0.65, PT_NOTE, 11, 1.4)
            self._text(slide, (SLIDE_W_IN - 9.0) / 2, 4.85, 9.0, 0.65,
                       supplement, pt=spt, color=TEXT_MUTED,
                       align=PP_ALIGN.CENTER, ls=1.4)

    # =====================================================
    # 5. numbered_cards — 番号丸バッジ付き横長カード（主力パターン）
    # =====================================================
    def numbered_cards(self, title, lead, items, takeaway=None,
                       accent_last=False):
        """
        items : list[(見出し, 本文)]（最大5件。本文は None 可）
        accent_last : True なら最終項目の番号を「+」のゴールドにする
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        items = items[:5]
        n = len(items)
        gap = 0.16
        h = (cb - ct - gap * (n - 1)) / n
        for i, (head, body) in enumerate(items):
            y = ct + i * (h + gap)
            self._card(slide, MX, y, CW, h)
            cy = y + h / 2
            is_plus = accent_last and i == n - 1
            d = min(0.52, h - 0.28)
            self._circle(slide, MX + 0.30 + (0.52 - d) / 2, cy - d / 2, d,
                         ACCENT if is_plus else PRIMARY,
                         glyph="+" if is_plus else str(i + 1),
                         glyph_pt=d * 36)
            tx, tw = MX + 1.08, CW - 1.40
            if body:
                bpt = _fit_pt(body, tw, h - 0.62, PT_NOTE, 10.5, 1.35)
                gh = 0.36 + _block_h(body, tw, bpt, 1.35)
                gy = cy - gh / 2
                self._text(slide, tx, gy, tw, 0.36, head, font=FONT_BOLD,
                           pt=_fit_one_line(head, tw, PT_BODY, 13),
                           color=TEXT, wrap=False)
                self._text(slide, tx, gy + 0.40, tw, gh - 0.36, body,
                           pt=bpt, color=TEXT_MUTED, ls=1.35)
            else:
                self._text(slide, tx, y, tw, h, head, font=FONT_BOLD,
                           pt=_fit_one_line(head, tw, PT_BODY, 13),
                           color=TEXT,
                           valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 6. icon_list — 小アイコンチップ付きリスト（icon-left-text-list）
    # =====================================================
    def icon_list(self, title, lead, items, takeaway=None):
        """
        items : list[(icon, 見出し, 本文)]（最大4件）
        icon  : PNGパス or 1〜2文字（例 "地"）。チップは56px相当に統一
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        items = items[:4]
        n = len(items)
        gap = 0.16
        h = (cb - ct - gap * (n - 1)) / n
        for i, (icon, head, body) in enumerate(items):
            y = ct + i * (h + gap)
            self._card(slide, MX, y, CW, h)
            tint, color = CHIP_COLORS[i % 3]
            self._icon_chip(slide, MX + 0.28, y + h / 2 - 0.29, 0.58,
                            icon, tint, color)
            tx, tw = MX + 1.12, CW - 1.45
            if body:
                bpt = _fit_pt(body, tw, h - 0.62, PT_NOTE, 10.5, 1.35)
                gh = 0.36 + _block_h(body, tw, bpt, 1.35)
                gy = y + h / 2 - gh / 2
                self._text(slide, tx, gy, tw, 0.36, head, font=FONT_BOLD,
                           pt=_fit_one_line(head, tw, PT_BODY, 13),
                           color=TEXT, wrap=False)
                self._text(slide, tx, gy + 0.40, tw, gh - 0.36, body,
                           pt=bpt, color=TEXT_MUTED, ls=1.35)
            else:
                self._text(slide, tx, y, tw, h, head, font=FONT_BOLD,
                           pt=_fit_one_line(head, tw, PT_BODY, 13),
                           color=TEXT,
                           valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 7. two_column — 2カラム本文（two-col-text-body）
    # =====================================================
    def two_column(self, title, lead, left, right, takeaway=None):
        """
        left / right : (見出し, list[箇条書き])
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        w = (CW - 0.4) / 2
        h = cb - ct
        for idx, ((head, items), x) in enumerate(
                [(left, MX), (right, MX + w + 0.4)]):
            color = PRIMARY if idx == 0 else SECONDARY
            self._card(slide, x, ct, w, h)
            self._text(slide, x + 0.32, ct + 0.26, w - 0.64, 0.40, head,
                       font=FONT_BOLD,
                       pt=_fit_one_line(head, w - 0.64, PT_BODY, 14),
                       color=TEXT, wrap=False)
            self._rect(slide, x + 0.32, ct + 0.74, 0.55, 0.035, color,
                       radius_in=0.017)
            self._dot_rows(slide, x + 0.32, ct + 1.09, w - 0.64, items,
                           pt=PT_BODY, dot_color=color, max_y=ct + h - 0.2,
                           spread_h=h - 1.49)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 8. before_after / problem_solution — 対比（中央シェブロン）
    # =====================================================
    def _compare(self, title, lead, left_label, left_items, right_label,
                 right_items, takeaway, left_mark, right_mark):
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        w = (CW - 1.05) / 2
        h = cb - ct
        rx = MX + w + 1.05
        self._chevron(slide, MX + w + (1.05 - 0.5) / 2, ct + h / 2 - 0.31,
                      w=0.5, h=0.62, color=PRIMARY)
        for x, label, items, mark, chip_bg, chip_fg, mk_color in [
            (MX, left_label, left_items, left_mark, GRAY_CHIP, TEXT_MUTED,
             TEXT_SUB),
            (rx, right_label, right_items, right_mark, PRIMARY, WHITE,
             PRIMARY),
        ]:
            self._card(slide, x, ct, w, h)
            lw = max(1.3, _est_w_in(label, PT_NOTE) + 0.55)
            self._rect(slide, x + 0.30, ct + 0.28, lw, 0.42, chip_bg,
                       radius_in=0.08)
            self._text(slide, x + 0.30, ct + 0.28, lw, 0.42, label,
                       font=FONT_BOLD, pt=PT_NOTE, color=chip_fg,
                       align=PP_ALIGN.CENTER,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
            iy = ct + 1.07
            row_gap = 0.18
            if len(items) > 1:                  # 縦に均等配分
                total = sum(max(_block_h(it, w - 0.98, PT_BODY, 1.3), 0.28)
                            for it in items)
                row_gap = min(0.55, max(0.18, (h - 1.47 - total)
                                        / (len(items) - 1)))
            for item in items:
                pt_i = _fit_pt(item, w - 0.98, ct + h - 0.25 - iy, PT_BODY,
                               10.5, 1.3)
                bh = _block_h(item, w - 0.98, pt_i, 1.3)
                self._circle(slide, x + 0.32, iy + max(bh, 0.24) / 2 - 0.11,
                             0.24, mk_color, glyph=mark, glyph_pt=10,
                             glyph_font=FONT_BOLD)
                self._text(slide, x + 0.68, iy, w - 0.98, bh + 0.1, item,
                           pt=pt_i, color=TEXT, ls=1.3)
                iy += max(bh, 0.28) + row_gap
        self._takeaway(slide, takeaway)

    def before_after(self, title, lead, before_items, after_items,
                     before_label="Before", after_label="After",
                     takeaway=None):
        """現状の課題（×）→ 改善後（✓）の対比。"""
        self._compare(title, lead, before_label, before_items,
                      after_label, after_items, takeaway, "✕", "✓")

    def problem_solution(self, title, lead, problem_items, solution_items,
                         problem_label="課題", solution_label="解決策",
                         takeaway=None):
        """課題（×）→ 解決策（✓）の対比。"""
        self._compare(title, lead, problem_label, problem_items,
                      solution_label, solution_items, takeaway, "✕", "✓")

    # =====================================================
    # 9. steps_horizontal — 横フロー（four-step-flow）
    # =====================================================
    def steps_horizontal(self, title, lead, steps, takeaway=None):
        """
        steps : list[(ラベル, 本文)]（3〜5件）
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        steps = steps[:5]
        n = len(steps)
        aw = 0.60                       # 矢印ゾーン幅
        w = (CW - (n - 1) * aw) / n
        label_only = all(not b for _, b in steps)
        h = min(cb - ct, 2.05 if label_only else 3.45)   # 縦長防止
        cy0 = ct + (cb - ct - h) / 2
        for i, (label, body) in enumerate(steps):
            x = MX + i * (w + aw)
            self._card(slide, x, cy0, w, h)
            self._circle(slide, x + w / 2 - 0.24, cy0 + 0.30, 0.48, PRIMARY,
                         glyph=str(i + 1), glyph_pt=17)
            # ラベルは1行優先で縮小（単語の途中改行を防ぐ）
            lpt = _fit_one_line(label, w - 0.30, PT_BODY, 11)
            self._text(slide, x + 0.15, cy0 + 0.95, w - 0.30, 0.62, label,
                       font=FONT_BOLD, pt=lpt,
                       color=TEXT, align=PP_ALIGN.CENTER, ls=1.15)
            if body:
                bpt = _fit_pt(body, w - 0.36, h - 1.85, PT_NOTE, 10, 1.35)
                self._text(slide, x + 0.18, cy0 + 1.68, w - 0.36, h - 1.85,
                           body, pt=bpt, color=TEXT_MUTED,
                           align=PP_ALIGN.CENTER, ls=1.35)
            if i < n - 1:
                self._chevron(slide, x + w + (aw - 0.36) / 2,
                              cy0 + h / 2 - 0.24, w=0.36, h=0.48)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 10. steps_vertical — 縦フロー（vertical-step-flow・逆算スケジュール）
    # =====================================================
    def steps_vertical(self, title, lead, steps, takeaway=None,
                       accent_last=True):
        """
        steps : list[(ラベル, 本文)]（3〜6件）。ラベル例 "STEP 1"・"10月"
        本文は str または (見出し, 補足) のタプル
        accent_last : 最終ステップ（到達点）をゴールドで強調
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        steps = steps[:6]
        n = len(steps)
        gap = 0.26
        h = (cb - ct - gap * (n - 1)) / n
        lw = 1.55
        for i, (label, body) in enumerate(steps):
            y = ct + i * (h + gap)
            is_last = accent_last and i == n - 1
            self._rect(slide, MX, y, lw, h, ACCENT if is_last else PRIMARY,
                       radius_in=0.10, shadow=True)
            self._text(slide, MX, y, lw, h, label, font=FONT_BOLD,
                       pt=_fit_label(label, lw - 0.2, h - 0.2, PT_BODY, 11),
                       color=WHITE, align=PP_ALIGN.CENTER,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, ls=1.15)
            cx = MX + lw + 0.30
            cw = CW - lw - 0.30
            self._card(slide, cx, y, cw, h)
            if isinstance(body, tuple):
                head, sub = body
                self._text(slide, cx + 0.32, y + h / 2 - 0.34, cw - 0.64,
                           0.36, head, font=FONT_BOLD,
                           pt=_fit_one_line(head, cw - 0.64, PT_BODY, 12),
                           color=TEXT, wrap=False)
                self._text(slide, cx + 0.32, y + h / 2 + 0.06, cw - 0.64,
                           h / 2 - 0.1, sub,
                           pt=_fit_one_line(sub, cw - 0.64, PT_NOTE, 10),
                           color=TEXT_MUTED, wrap=False)
            else:
                self._text(slide, cx + 0.32, y, cw - 0.64, h, body,
                           pt=_fit_pt(body, cw - 0.64, h - 0.2, PT_BODY,
                                      11, 1.3),
                           color=TEXT, valign=MSO_VERTICAL_ANCHOR.MIDDLE,
                           ls=1.3)
            if i < n - 1:
                self._down_arrow(slide, MX + lw / 2 - 0.17,
                                 y + h + (gap - 0.22) / 2, h=0.22)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 11. cards_3 — 3カラムアイコンカード（three-column-icon-card）
    # =====================================================
    def cards_3(self, title, lead, cards, takeaway=None):
        """
        cards : list[(icon, 見出し, 本文)] 3件
        icon  : PNGパス or 1〜2文字。チップ色は Primary/Secondary/Accent を循環
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        w = (CW - 0.48) / 3
        # 本文量に応じて高さを決める（縦長になりすぎない）
        max_body = max(_block_h(c[2], w - 0.56, PT_NOTE, 1.5)
                       for c in cards[:3])
        h = min(cb - ct, max(2.6, 1.80 + max_body + 0.35))
        cy0 = ct + (cb - ct - h) / 2
        for i, (icon, head, body) in enumerate(cards[:3]):
            x = MX + i * (w + 0.24)
            self._card(slide, x, cy0, w, h)
            tint, color = CHIP_COLORS[i % 3]
            self._icon_chip(slide, x + 0.28, cy0 + 0.28, 0.58, icon, tint,
                            color)
            self._text(slide, x + 0.28, cy0 + 1.04, w - 0.56, 0.40, head,
                       font=FONT_BOLD,
                       pt=_fit_one_line(head, w - 0.56, PT_BODY, 12.5),
                       color=TEXT, wrap=False)
            self._text(slide, x + 0.28, cy0 + 1.55, w - 0.56, h - 1.80, body,
                       pt=_fit_pt(body, w - 0.56, h - 1.80, PT_NOTE, 10.5,
                                  1.5),
                       color=TEXT_MUTED, ls=1.5)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 12. cards_2x2 — 4カード 2×2 グリッド（four-card-2x2）
    # =====================================================
    def cards_2x2(self, title, lead, cards, takeaway=None):
        """
        cards : list[(見出し, 本文)] 4件
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        gap = 0.22
        w = (CW - gap) / 2
        h = (cb - ct - gap) / 2
        for i, (head, body) in enumerate(cards[:4]):
            x = MX + (i % 2) * (w + gap)
            y = ct + (i // 2) * (h + gap)
            self._card(slide, x, y, w, h)
            self._circle(slide, x + 0.28, y + 0.26, 0.36, PRIMARY,
                         glyph=str(i + 1), glyph_pt=13)
            self._text(slide, x + 0.80, y + 0.24, w - 1.05, 0.40, head,
                       font=FONT_BOLD,
                       pt=_fit_one_line(head, w - 1.05, PT_BODY, 12.5),
                       color=TEXT, valign=MSO_VERTICAL_ANCHOR.MIDDLE,
                       wrap=False)
            self._text(slide, x + 0.30, y + 0.86, w - 0.60, h - 1.09, body,
                       pt=_fit_pt(body, w - 0.60, h - 1.09, PT_NOTE, 10.5,
                                  1.4),
                       color=TEXT_MUTED, ls=1.4)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 13. figure_text — 図＋箇条書き（image-left-text-right）
    # =====================================================
    def figure_text(self, title, lead, image_path, points, caption=None,
                    takeaway=None, image_right=False, image_frac=0.44):
        """
        points : list[str]（ドット箇条書き）または list[(見出し, 本文)]
        image_frac : 画像カードの幅比（0.35〜0.55 推奨）
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        h = cb - ct
        img_w = CW * image_frac
        txt_w = CW - img_w - 0.40
        if image_right:
            tx, ix = MX, MX + txt_w + 0.40
        else:
            ix, tx = MX, MX + img_w + 0.40
        # 画像カード
        self._card(slide, ix, ct, img_w, h)
        cap_h = 0.36 if caption else 0.0
        if image_path:
            self._picture_fit(slide, image_path, ix + 0.18, ct + 0.18,
                              img_w - 0.36, h - 0.36 - cap_h)
        if caption:
            self._text(slide, ix + 0.18, ct + h - 0.44, img_w - 0.36, 0.32,
                       caption, font=FONT_LIGHT,
                       pt=_fit_one_line(caption, img_w - 0.36, PT_NOTE, 9),
                       color=TEXT_MUTED,
                       align=PP_ALIGN.CENTER, wrap=False)
        # テキスト側
        if points and isinstance(points[0], tuple):
            py = ct + 0.10
            for head, body in points:
                self._circle(slide, tx, py + 0.08, 0.11, PRIMARY)
                self._text(slide, tx + 0.28, py - 0.06, txt_w - 0.28, 0.36,
                           head, font=FONT_BOLD,
                           pt=_fit_one_line(head, txt_w - 0.28, PT_BODY, 13),
                           color=TEXT, wrap=False)
                bpt = _fit_pt(body, txt_w - 0.28, 1.2, PT_NOTE, 10.5, 1.35)
                self._text(slide, tx + 0.28, py + 0.34, txt_w - 0.28, 1.2,
                           body, pt=bpt, color=TEXT_MUTED, ls=1.35)
                py += 0.40 + _block_h(body, txt_w - 0.28, bpt, 1.35) + 0.24
        else:
            self._dot_rows(slide, tx, ct + 0.15, txt_w, points, pt=PT_BODY,
                           gap=0.22, max_y=cb - 0.1)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 14. table_compare — 比較表（comparison-matrix-table）
    # =====================================================
    def table_compare(self, title, lead, col_headers, rows,
                      highlight_col=None, takeaway=None, label_w=2.8):
        """
        col_headers   : list[str]（先頭は行ラベル列の見出し。空文字可）
        rows          : list[list[str]]（各行 = [行ラベル, 値1, 値2, ...]）
        highlight_col : 強調する列番号（1始まり・値列）
        セル値 ✓/○ は Primary、✕ はグレー、△ はゴールドで自動着色
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        n_cols = len(col_headers)
        val_w = (CW - label_w) / (n_cols - 1)
        head_h = 0.52
        row_h = (cb - ct - head_h) / len(rows)
        glyph_color = {"✓": PRIMARY, "○": PRIMARY, "◎": PRIMARY,
                       "✕": TEXT_SUB, "×": TEXT_SUB, "△": ACCENT_DARK}

        def cell_x(c):
            return MX + (label_w + (c - 1) * val_w if c > 0 else 0)

        def cell_w(c):
            return val_w if c > 0 else label_w

        for c, htxt in enumerate(col_headers):
            x = cell_x(c)
            fill = PRIMARY if c > 0 else SECONDARY
            if highlight_col and c == highlight_col:
                fill = ACCENT
            self._rect(slide, x, ct, cell_w(c), head_h, fill, line=SURFACE,
                       line_w=1.0)
            self._text(slide, x, ct, cell_w(c), head_h, htxt,
                       font=FONT_BOLD,
                       pt=_fit_one_line(htxt, cell_w(c) - 0.2, PT_NOTE, 11),
                       color=WHITE, align=PP_ALIGN.CENTER,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
        for r, row in enumerate(rows):
            y = ct + head_h + r * row_h
            for c in range(n_cols):
                x = cell_x(c)
                fill = SURFACE
                if highlight_col and c == highlight_col:
                    fill = ACCENT_TINT
                self._rect(slide, x, y, cell_w(c), row_h, fill, line=BORDER,
                           line_w=0.75)
                val = row[c] if c < len(row) else ""
                if c == 0:
                    self._text(slide, x + 0.20, y, cell_w(c) - 0.35, row_h,
                               val, font=FONT_BOLD,
                               pt=_fit_label(val, cell_w(c) - 0.35,
                                             row_h - 0.1, PT_BODY, 10.5,
                                             ls=1.2, one_line_min=12),
                               color=TEXT, valign=MSO_VERTICAL_ANCHOR.MIDDLE,
                               ls=1.2)
                else:
                    col = glyph_color.get(val.strip(), TEXT)
                    fnt = FONT_BOLD if val.strip() in glyph_color else FONT_BODY
                    self._text(slide, x + 0.12, y, cell_w(c) - 0.24, row_h,
                               val, font=fnt,
                               pt=_fit_pt(val, cell_w(c) - 0.24, row_h - 0.1,
                                          PT_BODY, 10, 1.2),
                               color=col, align=PP_ALIGN.CENTER,
                               valign=MSO_VERTICAL_ANCHOR.MIDDLE, ls=1.2)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 15. kpi_3 — 3つの大数字（three-kpi-big-number）
    # =====================================================
    def kpi_3(self, title, lead, kpis, takeaway=None):
        """
        kpis : list[(値, 指標名, 補足)] 3件。値の例 "20倍" "±1桁" "5誌"
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        w = (CW - 0.48) / 3
        h = cb - ct
        for i, (value, label, note) in enumerate(kpis[:3]):
            x = MX + i * (w + 0.24)
            self._card(slide, x, ct, w, h)
            self._text(slide, x + 0.2, ct + h * 0.16, w - 0.4, 0.34, label,
                       font=FONT_BOLD, pt=14, color=TEXT_MUTED,
                       align=PP_ALIGN.CENTER, wrap=False)
            self._text(slide, x + 0.2, ct + h * 0.34, w - 0.4, 0.95, value,
                       font=FONT_HEAVY,
                       pt=_fit_pt(value, w - 0.4, 0.95, 40, 26, 1.0),
                       color=ACCENT_DARK, align=PP_ALIGN.CENTER,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
            self._rect(slide, x + w / 2 - 0.3, ct + h * 0.62, 0.6, 0.03,
                       TEAL_LIGHT, radius_in=0.015)
            if note:
                self._text(slide, x + 0.28, ct + h * 0.70, w - 0.56,
                           h * 0.26, note,
                           pt=_fit_pt(note, w - 0.56, h * 0.26, PT_NOTE, 9.5,
                                      1.35),
                           color=TEXT_MUTED, align=PP_ALIGN.CENTER, ls=1.35)
        self._takeaway(slide, takeaway)

    # =====================================================
    # 16. profile — 自己紹介（profile-bio）
    # =====================================================
    def profile(self, title, image_path=None, rows=None, keywords=None,
                lead=None, takeaway=None):
        """
        rows     : list[(ラベル, 値)]（所属・専門・経歴 など）
        keywords : list[str] — 下部に Primary ティントのタグとして並ぶ
        """
        slide = self._slide()
        ct = self._frame(slide, title, lead)
        cb = self._cb(takeaway)
        h = cb - ct
        rx, rw = MX, CW
        if image_path:
            img_w = 3.9
            self._card(slide, MX, ct, img_w, h)
            self._picture_fit(slide, image_path, MX + 0.22, ct + 0.22,
                              img_w - 0.44, h - 0.44)
            rx = MX + img_w + 0.45
            rw = CW - img_w - 0.45
        kw_h = 0.9 if keywords else 0.0
        rows_end = ct
        if rows:
            row_h = min(0.82, (h - kw_h) / len(rows))
            for i, (label, value) in enumerate(rows):
                y = ct + i * row_h
                self._text(slide, rx, y + 0.10, 1.9, 0.32, label,
                           font=FONT_BOLD,
                           pt=_fit_one_line(label, 1.9, PT_NOTE, 10),
                           color=TEXT_MUTED, wrap=False)
                self._text(slide, rx + 2.0, y + 0.06, rw - 2.0,
                           row_h - 0.12, value,
                           pt=_fit_pt(value, rw - 2.0, row_h - 0.14, PT_BODY,
                                      11, 1.25),
                           color=TEXT, ls=1.25)
                self._rect(slide, rx, y + row_h - 0.02, rw, 0.014, BORDER)
            rows_end = ct + len(rows) * row_h
        if keywords:
            ky = min(rows_end + 0.35, cb - 0.62)   # 行の直下に置く
            kx = rx
            for kw in keywords:
                kw_w = _est_w_in(kw, PT_NOTE) + 0.45
                if kx + kw_w > rx + rw:
                    break
                chip = self._rect(slide, kx, ky, kw_w, 0.42, PRIMARY_TINT,
                                  radius_in=0.21)
                chip.adjustments[0] = 0.5
                self._text(slide, kx, ky, kw_w, 0.42, kw, font=FONT_BOLD,
                           pt=PT_NOTE, color=PRIMARY, align=PP_ALIGN.CENTER,
                           valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)
                kx += kw_w + 0.18
        self._takeaway(slide, takeaway)

    # =====================================================
    # 17. closing — 締め（closing-slide）
    # =====================================================
    def closing(self, message, sub=None, contact=None):
        """
        message : お礼・締めの一言（中央大きく）
        contact : 連絡先。Primary のピルで表示
        """
        slide = self._slide()
        pt = _fit_pt(message, 10.6, 1.3, PT_KEYMSG, 19, 1.3)
        self._text(slide, (SLIDE_W_IN - 10.6) / 2, 2.35, 10.6, 1.3, message,
                   font=FONT_HEAVY, pt=pt, color=TEXT, align=PP_ALIGN.CENTER,
                   valign=MSO_VERTICAL_ANCHOR.MIDDLE, ls=1.35)
        if sub:
            self._text(slide, (SLIDE_W_IN - 9.6) / 2, 3.85, 9.6, 0.55, sub,
                       pt=_fit_pt(sub, 9.6, 0.55, PT_NOTE, 11, 1.4),
                       color=TEXT_SUB, align=PP_ALIGN.CENTER, ls=1.4)
        if contact:
            cpt = _fit_one_line(contact, 9.0, PT_NOTE, 11)
            w = min(9.5, _est_w_in(contact, cpt) + 1.5)
            x = (SLIDE_W_IN - w) / 2
            pill = self._rect(slide, x, 4.75, w, 0.62, PRIMARY,
                              radius_in=0.31, shadow=True)
            pill.adjustments[0] = 0.5
            self._text(slide, x, 4.75, w, 0.62, contact, font=FONT_BOLD,
                       pt=cpt, color=WHITE, align=PP_ALIGN.CENTER,
                       valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False)

    # =====================================================
    # Save
    # =====================================================
    def save(self, filename="output.pptx"):
        total = self._page
        for run, no in self._page_runs:
            run.text = f"{no} / {total}"
        self.prs.save(filename)
        return filename
